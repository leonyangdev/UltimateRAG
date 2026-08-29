"""对已启动的 UltimateRAG 执行一次可回收的 V2 全格式真实闭环验收。

脚本通过公开 API 动态生成并上传 Markdown、HTML、DOCX、XLSX、PPTX、PNG 和扫描 PDF，验证
Parser、READY、来源定位、Milvus 检索与百炼流式问答，最后清理临时知识库。

执行会真实调用百炼 Embedding、OCR 和 LLM，并产生少量模型用量：
``uv run python scripts/smoke_v2.py --api-url http://localhost:8000``
"""

import argparse
import json
import time
from collections.abc import Iterator
from dataclasses import dataclass
from io import BytesIO
from typing import Any
from uuid import uuid4

import httpx
from docx import Document as WordDocument
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


@dataclass(frozen=True, slots=True)
class SmokeDocument:
    """一份待上传的内存文档及其期望 Parser。"""

    filename: str
    mime_type: str
    content: bytes
    parser_name: str


def parse_arguments() -> argparse.Namespace:
    """解析 API 地址和验证问题。"""

    parser = argparse.ArgumentParser(description="Run the UltimateRAG V2 smoke test")
    parser.add_argument("--api-url", default="http://localhost:8000")
    parser.add_argument("--question", default="What is the UltimateRAG V2 smoke marker?")
    return parser.parse_args()


def require_success(response: httpx.Response, operation: str) -> None:
    """把非成功响应转换为带有限正文的验收错误。"""

    if response.is_success:
        return
    raise RuntimeError(f"{operation} failed: HTTP {response.status_code}: {response.text[:1000]}")


def ui_stream_events(response: httpx.Response) -> Iterator[dict[str, Any]]:
    """从 AI SDK SSE 中解析 JSON Event。"""

    for line in response.iter_lines():
        if not line or not line.startswith("data: "):
            continue
        payload = line.removeprefix("data: ")
        if payload == "[DONE]":
            return
        event = json.loads(payload)
        if not isinstance(event, dict):
            raise RuntimeError("stream event must be a JSON object")
        yield event


def build_documents() -> list[SmokeDocument]:
    """在内存中生成全部 V2 格式，避免仓库长期保存二进制 Fixture。"""

    marker = "UltimateRAG V2 smoke marker is 8291."
    documents = [
        SmokeDocument(
            "smoke.md",
            "text/markdown",
            f"# V2 Smoke\n\n{marker}".encode(),
            "markdown",
        ),
        SmokeDocument(
            "smoke.html",
            "text/html",
            f"<h1>V2 Smoke</h1><p>{marker}</p>".encode(),
            "html",
        ),
    ]

    word = WordDocument()
    word.add_heading("V2 Smoke", level=1)
    word.add_paragraph(marker)
    word_buffer = BytesIO()
    word.save(word_buffer)
    documents.append(
        SmokeDocument(
            "smoke.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            word_buffer.getvalue(),
            "docx",
        )
    )

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Smoke"
    sheet.append(["Name", "Value"])
    sheet.append(["UltimateRAG V2 smoke marker", 8291])
    excel_buffer = BytesIO()
    workbook.save(excel_buffer)
    documents.append(
        SmokeDocument(
            "smoke.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            excel_buffer.getvalue(),
            "xlsx",
        )
    )

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "V2 Smoke"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(1))
    text_box.text = marker
    ppt_buffer = BytesIO()
    presentation.save(ppt_buffer)
    documents.append(
        SmokeDocument(
            "smoke.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ppt_buffer.getvalue(),
            "pptx",
        )
    )

    # 大号高对比英文让 OCR Smoke Test 对本地字体环境不敏感；同一图像同时用于图片和扫描 PDF。
    image = Image.new("RGB", (1400, 320), "white")
    font = ImageFont.load_default(size=52)
    ImageDraw.Draw(image).text((40, 110), marker, fill="black", font=font)
    png_buffer = BytesIO()
    image.save(png_buffer, format="PNG")
    documents.append(SmokeDocument("smoke.png", "image/png", png_buffer.getvalue(), "image-ocr"))

    pdf_buffer = BytesIO()
    image.save(pdf_buffer, format="PDF", resolution=150)
    documents.append(
        SmokeDocument(
            "smoke.pdf",
            "application/pdf",
            pdf_buffer.getvalue(),
            "pdf-docling",
        )
    )
    return documents


def wait_until_ready(
    client: httpx.Client,
    document_id: str,
    filename: str,
    *,
    timeout_seconds: float = 600.0,
) -> dict[str, Any]:
    """轮询公开详情端点，验证异步任务最终进入 READY 或给出明确失败。"""

    deadline = time.monotonic() + timeout_seconds
    latest: dict[str, Any] = {}
    while time.monotonic() < deadline:
        response = client.get(f"/api/documents/{document_id}")
        require_success(response, f"poll {filename}")
        latest = response.json()
        if latest.get("status") == "READY":
            return latest
        if latest.get("status") == "FAILED":
            raise RuntimeError(f"{filename} background ingestion failed: {latest}")
        time.sleep(1)
    raise TimeoutError(f"{filename} did not finish background ingestion: {latest}")


def _verify_locator(parser_name: str, payload: dict[str, Any]) -> None:
    """验证格式特有 Locator 已从向量检索返回，而不只在 Parser 内短暂存在。"""

    locator = payload.get("locator")
    if not isinstance(locator, dict):
        raise RuntimeError(f"retrieval result has no locator: {payload}")
    if parser_name == "pdf-docling" and locator.get("page") is None:
        raise RuntimeError(f"PDF retrieval result has no page: {payload}")
    if parser_name == "xlsx" and not locator.get("sheet"):
        raise RuntimeError(f"XLSX retrieval result has no sheet: {payload}")
    if parser_name == "pptx" and locator.get("slide") is None:
        raise RuntimeError(f"PPTX retrieval result has no slide: {payload}")


def run_smoke_test(api_url: str, question: str) -> None:
    """执行 Create → Upload All → Retrieve → Stream Chat → Delete 完整验收。"""

    knowledge_base_id: str | None = None
    primary_error: Exception | None = None
    cleanup_error: str | None = None
    documents = build_documents()
    with httpx.Client(base_url=api_url.rstrip("/"), timeout=300.0, trust_env=False) as client:
        try:
            health = client.get("/api/health")
            require_success(health, "health check")
            create = client.post(
                "/api/knowledge-bases",
                json={
                    "name": f"V2 Smoke {uuid4().hex[:8]}",
                    "description": "由 scripts/smoke_v2.py 创建，可安全清理。",
                },
            )
            require_success(create, "create knowledge base")
            knowledge_base_id = str(create.json()["id"])
            print(f"[1/5] Knowledge base created: {knowledge_base_id}")

            accepted_documents: list[tuple[str, SmokeDocument]] = []
            for document in documents:
                response = client.post(
                    f"/api/knowledge-bases/{knowledge_base_id}/documents",
                    files={
                        "file": (document.filename, document.content, document.mime_type),
                    },
                )
                require_success(response, f"upload {document.filename}")
                payload = response.json()
                if response.status_code != 202 or payload.get("status") != "PENDING":
                    raise RuntimeError(
                        f"{document.filename} was not accepted asynchronously: "
                        f"HTTP {response.status_code} {payload}"
                    )
                accepted_documents.append((str(payload["id"]), document))

            uploaded: dict[str, str] = {}
            for document_id, document in accepted_documents:
                payload = wait_until_ready(client, document_id, document.filename)
                if payload.get("parser_name") != document.parser_name:
                    raise RuntimeError(
                        f"{document.filename} selected {payload.get('parser_name')}, "
                        f"expected {document.parser_name}"
                    )
                uploaded[document_id] = document.parser_name
            print(f"[2/5] All {len(documents)} V2 formats reached READY")

            retrieval = client.post(
                "/api/retrieval/search",
                json={"knowledge_base_id": knowledge_base_id, "query": question, "top_k": 20},
            )
            require_success(retrieval, "dense retrieval")
            results = retrieval.json()
            if not results:
                raise RuntimeError("retrieval returned no chunks")
            located_formats: set[str] = set()
            for result in results:
                parser_name = uploaded.get(str(result.get("document_id")))
                if parser_name in {"pdf-docling", "xlsx", "pptx"}:
                    _verify_locator(parser_name, result)
                    located_formats.add(parser_name)
            if located_formats != {"pdf-docling", "xlsx", "pptx"}:
                raise RuntimeError(
                    "top-20 retrieval did not return every locator format: "
                    f"found={sorted(located_formats)}"
                )
            print(f"[3/5] Retrieval returned {len(results)} chunks with V2 locators")

            with client.stream(
                "POST",
                "/api/chat/stream",
                json={"knowledge_base_id": knowledge_base_id, "question": question, "top_k": 5},
            ) as stream:
                require_success(stream, "stream chat")
                events = list(ui_stream_events(stream))
            event_types = [str(event.get("type")) for event in events]
            answer = "".join(
                str(event.get("delta", "")) for event in events if event.get("type") == "text-delta"
            )
            evidence = [event for event in events if event.get("type") == "data-retrieval"]
            if not answer.strip() or not evidence or "finish" not in event_types:
                raise RuntimeError(f"invalid chat stream: events={event_types}, answer={answer!r}")
            evidence_data = evidence[0].get("data")
            if not isinstance(evidence_data, dict) or not evidence_data.get("citations"):
                raise RuntimeError("stream retrieval data contains no citations")
            print(f"[4/5] Stream answer and citations verified ({len(answer)} chars)")
        except Exception as exc:
            primary_error = exc
        finally:
            if knowledge_base_id is not None:
                try:
                    cleanup = client.delete(f"/api/knowledge-bases/{knowledge_base_id}")
                except httpx.HTTPError as exc:
                    cleanup_error = f"knowledge_base_id={knowledge_base_id}, cleanup error: {exc}"
                else:
                    if cleanup.is_success:
                        print("[5/5] Temporary knowledge base deleted")
                    else:
                        cleanup_error = (
                            f"knowledge_base_id={knowledge_base_id}, cleanup HTTP "
                            f"{cleanup.status_code}: {cleanup.text[:500]}"
                        )

    if primary_error is not None:
        if cleanup_error:
            primary_error.add_note(f"Additional cleanup failure: {cleanup_error}")
        raise primary_error
    if cleanup_error:
        raise RuntimeError(cleanup_error)


def main() -> None:
    """运行命令行验收并输出明确发布信号。"""

    arguments = parse_arguments()
    run_smoke_test(arguments.api_url, arguments.question)
    print("UltimateRAG V2 smoke test passed.")


if __name__ == "__main__":
    main()
