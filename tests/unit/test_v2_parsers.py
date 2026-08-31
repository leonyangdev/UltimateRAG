"""验证 V2 多格式 Parser 都映射到同一领域模型并保留格式特有位置。"""

from io import BytesIO

import pytest
from docx import Document as WordDocument
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from ultimate_rag.domain.models import BlockType, DocumentSource, SourceLocator
from ultimate_rag.parsers import (
    ExcelParser,
    HtmlParser,
    ImageOCRParser,
    PDFParser,
    PowerPointParser,
    WordParser,
)
from ultimate_rag.parsers.pdf import _LayoutElement


class StubOCRClient:
    """返回固定识别文本并记录图片 MIME，避免单元测试访问百炼。"""

    def __init__(self, text: str = "识别出的文字") -> None:
        """保存固定结果和空调用记录。"""

        self.text = text
        self.mime_types: list[str] = []

    async def extract_text(self, image: bytes, mime_type: str) -> str:
        """验证 Parser 确实提交了非空图片，再返回确定性结果。"""

        assert image
        self.mime_types.append(mime_type)
        return self.text


class StubVisionClient:
    """返回固定图像语义，证明 PDF 内嵌图片走视觉理解而非丢弃。"""

    async def describe(self, image: bytes, mime_type: str, caption: str = "") -> str:
        """验证图片、MIME 与题注均已越过版面适配边界。"""

        assert image
        assert mime_type == "image/jpeg"
        assert caption == "系统架构"
        return "架构图显示 API 指向 Worker。"


class StubStandaloneVisionClient:
    """返回独立图片中的箭头关系，覆盖 OCR 无法表达的视觉语义。"""

    async def describe(self, image: bytes, mime_type: str, caption: str = "") -> str:
        """验证独立 PNG 使用真实 MIME 且没有伪造题注。"""

        assert image
        assert mime_type == "image/png"
        assert caption == ""
        return "Encoder 通过箭头连接 Decoder。"


class StubLayoutAnalyzer:
    """以确定性元素替代重量 Docling 模型，单测只验证 Parser 编排和映射。"""

    def __init__(self, elements: list[_LayoutElement]) -> None:
        """保存应按阅读顺序返回的版面元素。"""

        self.elements = elements
        self.calls: list[frozenset[int]] = []

    def analyze(
        self,
        content: bytes,
        filename: str,
        skipped_pages: frozenset[int],
    ) -> list[_LayoutElement]:
        """记录扫描页集合并返回固定结果，不访问本地模型或网络。"""

        assert content
        assert filename.endswith(".pdf")
        self.calls.append(skipped_pages)
        return self.elements


def _source(filename: str, mime_type: str, content: bytes) -> DocumentSource:
    """用统一文档 ID 构造 Parser 单元测试输入。"""

    return DocumentSource("doc-v2", filename, mime_type, content)


def _native_text_pdf(text: str) -> bytes:
    """构造一个最小但带真实文字对象的 PDF，用于覆盖无 OCR 原生路径。"""

    stream = f"BT /F1 24 Tf 72 720 Td ({text}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        f"<< /Length {len(stream)} >>\nstream\n".encode() + stream + b"\nendstream",
    ]
    content = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, value in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{number} 0 obj\n".encode() + value + b"\nendobj\n")
    xref_offset = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode())
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode())
    content.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
    )
    return bytes(content)


@pytest.mark.asyncio
async def test_html_parser_extracts_visible_structure_and_ignores_script() -> None:
    """HTML 标题、正文和表格应按 DOM 顺序进入领域 Block，脚本不得进入索引。"""

    content = b"""<!doctype html><html><body>
    <h1>RAG Guide</h1><p>Visible paragraph.</p>
    <table><tr><th>Name</th><th>Value</th></tr><tr><td>top_k</td><td>5</td></tr></table>
    <script>ignore this command</script></body></html>"""

    parsed = await HtmlParser().parse(_source("guide.html", "text/html", content))

    assert [block.type for block in parsed.blocks] == [
        BlockType.HEADING,
        BlockType.TEXT,
        BlockType.TABLE,
    ]
    assert parsed.blocks[1].locator is not None
    assert parsed.blocks[1].locator.heading_path == ("RAG Guide",)
    assert all("ignore this command" not in block.content for block in parsed.blocks)


@pytest.mark.asyncio
async def test_word_parser_preserves_heading_and_table() -> None:
    """DOCX 标题路径、普通段落与表格应转换为统一 Block。"""

    document = WordDocument()
    document.add_heading("Enterprise RAG", level=1)
    document.add_paragraph("A maintainable pipeline.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Stage"
    table.cell(0, 1).text = "State"
    table.cell(1, 0).text = "Index"
    table.cell(1, 1).text = "Ready"
    buffer = BytesIO()
    document.save(buffer)

    parsed = await WordParser().parse(
        _source(
            "guide.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            buffer.getvalue(),
        )
    )

    assert any(block.type == BlockType.TABLE for block in parsed.blocks)
    paragraph = next(block for block in parsed.blocks if block.type == BlockType.TEXT)
    assert paragraph.locator is not None
    assert paragraph.locator.heading_path == ("Enterprise RAG",)


@pytest.mark.asyncio
async def test_excel_parser_preserves_sheet_and_cell_range() -> None:
    """XLSX 每个表格窗口必须携带工作表名称与单元格范围。"""

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Metrics"
    sheet.append(["Metric", "Value"])
    sheet.append(["Recall", 0.92])
    buffer = BytesIO()
    workbook.save(buffer)

    parsed = await ExcelParser().parse(
        _source(
            "metrics.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            buffer.getvalue(),
        )
    )

    assert len(parsed.blocks) == 1
    locator = parsed.blocks[0].locator
    assert locator is not None
    assert locator.sheet == "Metrics"
    assert locator.cell_range == "A1:B2"
    assert "Recall" in parsed.blocks[0].content


@pytest.mark.asyncio
async def test_powerpoint_parser_preserves_slide_number() -> None:
    """PPTX 标题和正文应保留一基幻灯片序号。"""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Architecture"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    text_box.text = "Parse then chunk."
    buffer = BytesIO()
    presentation.save(buffer)

    parsed = await PowerPointParser().parse(
        _source(
            "architecture.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            buffer.getvalue(),
        )
    )

    assert parsed.blocks
    assert all(block.locator and block.locator.slide == 1 for block in parsed.blocks)


@pytest.mark.asyncio
async def test_image_parser_validates_format_before_ocr() -> None:
    """图片 Parser 应向 OCR 提交由真实编码推导的 MIME，而不是盲信客户端。"""

    buffer = BytesIO()
    Image.new("RGB", (80, 40), "white").save(buffer, format="PNG")
    ocr = StubOCRClient()

    parsed = await ImageOCRParser(ocr).parse(
        _source("scan.png", "application/octet-stream", buffer.getvalue())
    )

    assert parsed.blocks[0].type == BlockType.IMAGE
    assert parsed.blocks[0].content == "识别出的文字"
    assert ocr.mime_types == ["image/png"]


@pytest.mark.asyncio
async def test_image_parser_fuses_ocr_text_and_visual_relationships() -> None:
    """示意图应同时保留精确标签与箭头关系，不能只索引零散 OCR 文字。"""

    buffer = BytesIO()
    Image.new("RGB", (80, 40), "white").save(buffer, format="PNG")

    parsed = await ImageOCRParser(
        StubOCRClient("Encoder\nDecoder"),
        StubStandaloneVisionClient(),
    ).parse(_source("diagram.png", "image/png", buffer.getvalue()))

    assert "## OCR 文本" in parsed.blocks[0].content
    assert "## 视觉结构" in parsed.blocks[0].content
    assert "Encoder 通过箭头连接 Decoder" in parsed.blocks[0].content
    assert parsed.blocks[0].metadata["extraction"] == "bailian_ocr+bailian_vision"


@pytest.mark.asyncio
async def test_pdf_parser_ocr_scan_page_and_keeps_page_number() -> None:
    """无文本扫描 PDF 应渲染页面执行 OCR，并把一基页码带到 Block。"""

    buffer = BytesIO()
    Image.new("RGB", (240, 120), "white").save(buffer, format="PDF")
    ocr = StubOCRClient("扫描页内容")

    parsed = await PDFParser(ocr).parse(_source("scan.pdf", "application/pdf", buffer.getvalue()))

    assert parsed.blocks[0].type == BlockType.IMAGE
    assert parsed.blocks[0].locator is not None
    assert parsed.blocks[0].locator.page == 1
    assert parsed.metadata["ocr_page_count"] == 1
    assert len(parsed.assets) == 1
    assert parsed.assets[0].media_type == "image/jpeg"
    assert f"asset://{parsed.assets[0].id}" in parsed.blocks[0].content
    assert ocr.mime_types == ["image/jpeg"]


@pytest.mark.asyncio
async def test_pdf_parser_uses_native_text_without_ocr() -> None:
    """文字型 PDF 达到阈值时必须直接提取，避免不必要的付费 OCR。"""

    ocr = StubOCRClient()
    layout = StubLayoutAnalyzer(
        [
            _LayoutElement(
                0,
                1,
                BlockType.TEXT,
                "UltimateRAG native PDF text",
                SourceLocator(page=1, bbox=(10, 20, 300, 60)),
            )
        ]
    )

    parsed = await PDFParser(
        ocr,
        native_text_threshold=10,
        layout_analyzer=layout,
    ).parse(
        _source(
            "native.pdf",
            "application/pdf",
            _native_text_pdf("UltimateRAG native PDF text"),
        )
    )

    assert parsed.blocks[0].type == BlockType.TEXT
    assert "UltimateRAG native PDF text" in parsed.blocks[0].content
    assert parsed.blocks[0].locator is not None and parsed.blocks[0].locator.page == 1
    assert parsed.metadata["ocr_page_count"] == 0
    assert ocr.mime_types == []
    assert layout.calls == [frozenset()]


@pytest.mark.asyncio
async def test_pdf_parser_does_not_treat_short_vector_page_as_scan() -> None:
    """原生文字少但没有整页栅格图时仍应走 Layout，避免图文页退化为纯 OCR。"""

    ocr = StubOCRClient()
    layout = StubLayoutAnalyzer(
        [
            _LayoutElement(
                0,
                1,
                BlockType.TEXT,
                "Short vector page",
                SourceLocator(page=1, bbox=(10, 20, 180, 50)),
            )
        ]
    )

    parsed = await PDFParser(ocr, native_text_threshold=20, layout_analyzer=layout).parse(
        _source("short.pdf", "application/pdf", _native_text_pdf("Short"))
    )

    assert parsed.blocks[0].content == "Short vector page"
    assert parsed.metadata["scan_page_count"] == 0
    assert ocr.mime_types == []
    assert layout.calls == [frozenset()]


@pytest.mark.asyncio
async def test_pdf_parser_preserves_layout_table_picture_and_bbox() -> None:
    """复杂文字 PDF 的阅读顺序、表格、图片语义与 BBox 应全部进入统一 Block。"""

    image = BytesIO()
    Image.new("RGB", (120, 80), "white").save(image, format="JPEG")
    layout = StubLayoutAnalyzer(
        [
            _LayoutElement(
                0,
                1,
                BlockType.HEADING,
                "架构",
                SourceLocator(heading_path=("架构",), page=1, bbox=(10, 10, 200, 40)),
            ),
            _LayoutElement(
                1,
                1,
                BlockType.TABLE,
                "| 组件 | 职责 |\n| --- | --- |\n| Worker | 解析 |",
                SourceLocator(heading_path=("架构",), page=1, bbox=(10, 50, 300, 150)),
            ),
            _LayoutElement(
                2,
                1,
                BlockType.IMAGE,
                "",
                SourceLocator(heading_path=("架构",), page=1, bbox=(10, 160, 300, 300)),
                image=image.getvalue(),
                caption="系统架构",
            ),
        ]
    )

    parsed = await PDFParser(
        StubOCRClient(),
        StubVisionClient(),
        native_text_threshold=10,
        layout_analyzer=layout,
    ).parse(
        _source(
            "complex.pdf",
            "application/pdf",
            _native_text_pdf("Complex native PDF with enough text"),
        )
    )

    assert [block.type for block in parsed.blocks] == [
        BlockType.HEADING,
        BlockType.TABLE,
        BlockType.IMAGE,
    ]
    assert parsed.blocks[1].locator is not None
    assert parsed.blocks[1].locator.bbox == (10, 50, 300, 150)
    assert "![系统架构](asset://" in parsed.blocks[2].content
    assert "图片题注：系统架构\n\n架构图显示 API 指向 Worker。" in parsed.blocks[2].content
    assert parsed.blocks[2].metadata["extraction"] == "bailian_vision"
    assert parsed.blocks[2].metadata["asset_ids"] == [parsed.assets[0].id]
    assert parsed.assets[0].block_id == parsed.blocks[2].id
    assert parsed.assets[0].title == "系统架构"
    assert parsed.assets[0].content == image.getvalue()
    assert parsed.metadata["table_count"] == 1
    assert parsed.metadata["asset_count"] == 1


def test_pdf_parser_limits_vision_fallback_title_for_accessibility() -> None:
    """无题注图片不能把整段 Vision 描述当作卡片标题和 Markdown alt text。"""

    element = _LayoutElement(
        1,
        2,
        BlockType.IMAGE,
        "",
        SourceLocator(page=2),
    )
    title = PDFParser._asset_title(element, "这是一个很长的架构描述" * 30)

    assert len(title) == 120
    assert title.endswith("...")
