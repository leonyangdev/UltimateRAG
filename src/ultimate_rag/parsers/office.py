"""DOCX、XLSX 与 PPTX 到统一文档模型的 Office Parser。

模块职责：
    使用各格式成熟库提取标题、正文、列表与表格，立即映射为 UltimateRAG 的 Block 和
    SourceLocator。所有同步 ZIP/XML 解析在线程中执行，避免阻塞 FastAPI Event Loop。

安全边界：
    三种 OOXML 格式在交给第三方库前统一检查 ZIP 条目数、解压总量和异常压缩比，降低压缩炸弹
    风险。V2 只支持现代 OOXML，不支持带宏格式或旧版 `.doc`、`.xls`、`.ppt`。

已知限制：
    V2 提取 Office 原生文字与表格，不对文档内嵌图片再次 OCR，不计算 Word 真实页码，也不解析
    图表背后的数据模型；这些内容不会阻止同一文档中的可见文字进入 RAG。
"""

import asyncio
from io import BytesIO

from docx import Document as OpenDocument
from docx.table import Table as WordTable
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation

from ultimate_rag.domain.exceptions import InvalidDocumentError
from ultimate_rag.domain.models import (
    Block,
    BlockType,
    DocumentSource,
    ParsedDocument,
    SourceLocator,
)
from ultimate_rag.parsers._shared import (
    stable_block,
    supports_source,
    table_to_markdown,
    validate_ooxml_archive,
)


class WordParser:
    """把 DOCX 段落、标题、列表和表格转换为有序 Block。"""

    name = "docx"
    version = "2.0"
    _EXTENSIONS = frozenset({".docx"})
    _MIME_TYPES = frozenset(
        {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
    )

    def supports(self, source: DocumentSource) -> bool:
        """按 DOCX 扩展名和标准 MIME 识别来源。"""

        return supports_source(source, self._EXTENSIONS, self._MIME_TYPES)

    async def parse(self, source: DocumentSource) -> ParsedDocument:
        """在线程中解析 DOCX，保留标题路径和表格结构。"""

        validate_ooxml_archive(source.content)
        return await asyncio.to_thread(self._parse_sync, source)

    def _parse_sync(self, source: DocumentSource) -> ParsedDocument:
        """执行 python-docx 同步解析并映射领域模型。"""

        try:
            document = OpenDocument(BytesIO(source.content))
            blocks: list[Block] = []
            heading_path: list[str] = []
            for element in document.iter_inner_content():
                if isinstance(element, Paragraph):
                    content = element.text.strip()
                    if not content:
                        continue
                    style_name = element.style.name if element.style is not None else ""
                    block_type = BlockType.TEXT
                    if style_name.startswith("Heading"):
                        level_text = style_name.removeprefix("Heading").strip()
                        level = int(level_text) if level_text.isdigit() else 1
                        heading_path = heading_path[: level - 1] + [content]
                        block_type = BlockType.HEADING
                    elif style_name.startswith("List"):
                        block_type = BlockType.LIST
                    locator = SourceLocator(heading_path=tuple(heading_path))
                    blocks.append(
                        stable_block(source.document_id, len(blocks), block_type, content, locator)
                    )
                elif isinstance(element, WordTable):
                    content = table_to_markdown(
                        [[cell.text for cell in row.cells] for row in element.rows]
                    )
                    if content:
                        locator = SourceLocator(heading_path=tuple(heading_path))
                        blocks.append(
                            stable_block(
                                source.document_id,
                                len(blocks),
                                BlockType.TABLE,
                                content,
                                locator,
                            )
                        )
        except InvalidDocumentError:
            raise
        except Exception as exc:
            raise InvalidDocumentError("DOCX 文件损坏或无法解析") from exc
        return _parsed(source, self.name, self.version, blocks)


class ExcelParser:
    """把 XLSX 工作表按有限行窗口转换为 Markdown Table Block。"""

    name = "xlsx"
    version = "2.0"
    _EXTENSIONS = frozenset({".xlsx"})
    _MIME_TYPES = frozenset({"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"})
    _MAX_SHEETS = 100
    _MAX_CELLS_PER_SHEET = 200_000
    _ROWS_PER_BLOCK = 100

    def supports(self, source: DocumentSource) -> bool:
        """按 XLSX 扩展名和标准 MIME 识别来源。"""

        return supports_source(source, self._EXTENSIONS, self._MIME_TYPES)

    async def parse(self, source: DocumentSource) -> ParsedDocument:
        """在线程中以 read-only 模式读取工作表，避免复制完整 Workbook。"""

        validate_ooxml_archive(source.content)
        return await asyncio.to_thread(self._parse_sync, source)

    def _parse_sync(self, source: DocumentSource) -> ParsedDocument:
        """执行 openpyxl 同步读取，并给每个表格窗口记录 Sheet 与 Range。"""

        try:
            workbook = load_workbook(BytesIO(source.content), read_only=True, data_only=False)
            if len(workbook.worksheets) > self._MAX_SHEETS:
                raise InvalidDocumentError("XLSX 工作表数量超过限制")
            blocks: list[Block] = []
            for worksheet in workbook.worksheets:
                if worksheet.max_row * worksheet.max_column > self._MAX_CELLS_PER_SHEET:
                    raise InvalidDocumentError(f"XLSX 工作表 {worksheet.title} 单元格范围过大")
                rows = list(worksheet.iter_rows(values_only=True))
                for start in range(0, len(rows), self._ROWS_PER_BLOCK):
                    window = rows[start : start + self._ROWS_PER_BLOCK]
                    content = table_to_markdown(window)
                    if not content:
                        continue
                    end_row = min(start + len(window), worksheet.max_row)
                    end_cell = worksheet.cell(end_row, worksheet.max_column).coordinate
                    cell_range = f"A{start + 1}:{end_cell}"
                    locator = SourceLocator(sheet=worksheet.title, cell_range=cell_range)
                    blocks.append(
                        stable_block(
                            source.document_id,
                            len(blocks),
                            BlockType.TABLE,
                            content,
                            locator,
                        )
                    )
            workbook.close()
        except InvalidDocumentError:
            raise
        except Exception as exc:
            raise InvalidDocumentError("XLSX 文件损坏或无法解析") from exc
        return _parsed(source, self.name, self.version, blocks)


class PowerPointParser:
    """把 PPTX 每张幻灯片中的标题、文本框和表格转换为 Block。"""

    name = "pptx"
    version = "2.0"
    _EXTENSIONS = frozenset({".pptx"})
    _MIME_TYPES = frozenset(
        {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    )

    def supports(self, source: DocumentSource) -> bool:
        """按 PPTX 扩展名和标准 MIME 识别来源。"""

        return supports_source(source, self._EXTENSIONS, self._MIME_TYPES)

    async def parse(self, source: DocumentSource) -> ParsedDocument:
        """在线程中解析 PPTX，并记录一基幻灯片序号。"""

        validate_ooxml_archive(source.content)
        return await asyncio.to_thread(self._parse_sync, source)

    def _parse_sync(self, source: DocumentSource) -> ParsedDocument:
        """执行 python-pptx 同步解析并保持 Shape 原始顺序。"""

        try:
            presentation = Presentation(BytesIO(source.content))
            blocks: list[Block] = []
            for slide_number, slide in enumerate(presentation.slides, start=1):
                title_shape = slide.shapes.title
                title = title_shape.text.strip() if title_shape is not None else ""
                heading_path = (title,) if title else ()
                if title:
                    blocks.append(
                        stable_block(
                            source.document_id,
                            len(blocks),
                            BlockType.HEADING,
                            title,
                            SourceLocator(heading_path=heading_path, slide=slide_number),
                        )
                    )
                for shape in slide.shapes:
                    if shape is title_shape:
                        continue
                    locator = SourceLocator(heading_path=heading_path, slide=slide_number)
                    if shape.has_table:
                        content = table_to_markdown(
                            [[cell.text for cell in row.cells] for row in shape.table.rows]
                        )
                        block_type = BlockType.TABLE
                    elif shape.has_text_frame:
                        content = shape.text.strip()
                        block_type = BlockType.TEXT
                    else:
                        continue
                    if content:
                        blocks.append(
                            stable_block(
                                source.document_id,
                                len(blocks),
                                block_type,
                                content,
                                locator,
                            )
                        )
        except Exception as exc:
            raise InvalidDocumentError("PPTX 文件损坏或无法解析") from exc
        return _parsed(source, self.name, self.version, blocks)


def _parsed(
    source: DocumentSource,
    parser_name: str,
    parser_version: str,
    blocks: list[Block],
) -> ParsedDocument:
    """统一拒绝无文本 Office 文件并构造 Parser 元数据。"""

    if not blocks:
        raise InvalidDocumentError(f"{source.filename} 没有可索引的文本或表格")
    return ParsedDocument(
        document_id=source.document_id,
        blocks=tuple(blocks),
        metadata={"parser": parser_name, "parser_version": parser_version},
    )
